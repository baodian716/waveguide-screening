"""Build the interview slide deck for waveguide-screening.

Run:
    py build_slides.py
Output:
    waveguide-screening.pptx (in this folder)
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

HERE = Path(__file__).parent
PROJECT_ROOT = HERE.parent
OUT_PATH = HERE / "waveguide-screening.pptx"

ACCENT = RGBColor(0x1F, 0x4E, 0x79)
GREY = RGBColor(0x55, 0x55, 0x55)


def add_title_only_slide(prs, title_text):
    layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(layout)
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = ACCENT
    return slide


def add_bullets(slide, bullets, left=0.7, top=1.6, width=8.6, height=5.0, font_size=20):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.word_wrap = True
    for i, (text, level) in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = level
        for run in p.runs:
            run.font.size = Pt(font_size if level == 0 else font_size - 4)
            if level == 0:
                run.font.bold = True
        p.space_after = Pt(8)


def add_subtitle(slide, text, top=1.05):
    box = slide.shapes.add_textbox(Inches(0.7), Inches(top), Inches(8.6), Inches(0.45))
    tf = box.text_frame
    tf.text = text
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = GREY
    tf.paragraphs[0].font.italic = True


def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ----- 1. Title -----
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(2.4), Inches(8.6), Inches(1.2))
    tf = title_box.text_frame
    tf.text = "waveguide-screening"
    tf.paragraphs[0].font.size = Pt(44)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = ACCENT

    sub = slide.shapes.add_textbox(Inches(0.7), Inches(3.6), Inches(8.6), Inches(0.8))
    sub.text_frame.text = (
        "Physics-guided parameter screening, refactored from a master's thesis"
    )
    sub.text_frame.paragraphs[0].font.size = Pt(20)
    sub.text_frame.paragraphs[0].font.color.rgb = GREY

    name = slide.shapes.add_textbox(Inches(0.7), Inches(6.4), Inches(8.6), Inches(0.5))
    name.text_frame.text = "[Your name]   |   [Date]"
    name.text_frame.paragraphs[0].font.size = Pt(14)
    name.text_frame.paragraphs[0].font.color.rgb = GREY

    # ----- 2. Problem -----
    slide = add_title_only_slide(prs, "Problem")
    add_subtitle(slide, "從碩論研究衍生的篩選需求")
    add_bullets(slide, [
        ("研究目標：找適合長距離傳輸的波導結構參數", 0),
        ("Search space：81 widths × 800 wavelengths = 64,800 candidates", 1),
        ("純暴力掃描無法直接挑出最佳組合", 1),
        ("需要的是 physics-guided screening：用物理條件先過濾、再排序", 0),
        ("挑戰：條件設計、流程組裝、結果視覺化", 1),
    ])

    # ----- 3. Method -----
    slide = add_title_only_slide(prs, "Method: three physics conditions + ranking")
    add_subtitle(slide, "用物理意義可解釋的條件做正交篩選")
    add_bullets(slide, [
        ("C1  Peak intensity at center axis  (top 5%)", 0),
        ("中心軸能量夠強 → 光仍集中在傳輸方向", 1),
        ("C2  Er / Ephi overlap integral  η ≥ 0.95", 0),
        ("場分布對稱 → 模態純度高", 1),
        ("C3  Half-divergence angle  < 1.5°", 0),
        ("光束發散夠低 → 能維持長距離傳輸", 1),
        ("→ 通過三條件後，依 transmission efficiency 排序取 Top-10", 0),
    ])

    # ----- 4. Architecture -----
    slide = add_title_only_slide(prs, "From research script to modular pipeline")
    add_subtitle(slide, "把 281 行 MATLAB 單檔重構成 6 個 Python 模組")
    add_bullets(slide, [
        ("原始狀態：MATLAB 單檔，I/O + 計算 + 篩選 + 視覺化全混", 0),
        ("重構後：模組化、可單元測試、易擴充", 0),
        ("config.py — 網格與閾值常數", 1),
        ("generate_demo_data.py — 合成資料生成", 1),
        ("metrics.py — 物理指標純函數", 1),
        ("screening.py — 條件套用 + 排序", 1),
        ("visualization.py — 圖表輸出", 1),
        ("main.py — pipeline 串接", 1),
    ])

    # ----- 5. Demo data strategy -----
    slide = add_title_only_slide(prs, "Demo data strategy")
    add_subtitle(slide, "為什麼公開版用 synthetic data")
    add_bullets(slide, [
        ("原始 simulation 資料來自專用光學/電磁模擬工具，未公開", 0),
        ("用 physics-inspired synthesis 替代", 0),
        ("Gaussian envelope × radial ringing", 1),
        ("在參數空間中設計 sweet spot，部分組合自然通過篩選", 1),
        ("驗證 pipeline 能正確識別這些 sweet spot", 1),
        ("8,200 組合 → AND 後 326 → Top-10", 0),
        ("Efficiency 12.8% ~ 13.7%，與原研究 5–23% 同數量級", 1),
    ])

    # ----- 6. Results -----
    slide = add_title_only_slide(prs, "Results")
    add_subtitle(slide, "通過三條件的 candidate 與 Top-10 排序結果")

    img1 = PROJECT_ROOT / "outputs" / "screening_and.png"
    img2 = PROJECT_ROOT / "outputs" / "top10_scatter.png"
    if img1.exists():
        slide.shapes.add_picture(
            str(img1), Inches(0.3), Inches(1.7), width=Inches(4.7)
        )
    if img2.exists():
        slide.shapes.add_picture(
            str(img2), Inches(5.0), Inches(1.7), width=Inches(4.7)
        )

    cap = slide.shapes.add_textbox(Inches(0.3), Inches(6.5), Inches(9.4), Inches(0.4))
    cap.text_frame.text = (
        "左：通過三條件的 326 個 candidate    "
        "右：Top-10（efficiency 0.128 ~ 0.137）"
    )
    cap.text_frame.paragraphs[0].font.size = Pt(14)
    cap.text_frame.paragraphs[0].font.color.rgb = GREY

    # ----- 7. Transferable skills -----
    slide = add_title_only_slide(prs, "What this demonstrates")
    add_subtitle(slide, "從這份專案能看出的可遷移能力")
    add_bullets(slide, [
        ("Metric design — 把領域知識轉成可量化、可組合的條件", 0),
        ("這個能力直接對應到 ML evaluation pipeline 的設計", 1),
        ("Pipeline thinking — 從研究腳本到 modular software", 0),
        ("Domain reasoning + Python tooling", 0),
        ("適合領域 AI 場景（材料、半導體、光學的 surrogate model 等）", 1),
        ("Honest framing", 0),
        ("這不是 AI 專案，是工程基本功展示。AI 作品另案進行。", 1),
    ])

    prs.save(OUT_PATH)
    print(f"Slide deck saved to {OUT_PATH}")


if __name__ == "__main__":
    build()
